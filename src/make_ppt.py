import os, cv2, re
from pptx.util import Inches,Cm
from config import *
from pptx import Presentation
from tqdm import tqdm
from tools import clean_images_in_path

def main():
    # 读取所有路径并排序
    image_paths = [os.path.join(output_picture_path, f) for f in os.listdir(output_picture_path) if f.endswith('.png')]
    image_paths.sort(
        key=lambda x: int(re.search(r'frame_(\d+)', x).group(1)) if re.search(r'frame_(\d+)', x) else 0)

    if len(image_paths)==0:
        print("未读取到图片")
        return 0

    # 创建一个新的PPT对象
    prs = Presentation()
    # 设置幻灯片的页面大小和长宽比
    img1 = cv2.imread(image_paths[0])
    img1_aspect_radio = img1.shape[1] / img1.shape[0]
    width = Cm(25)  # 幻灯片宽度（厘米）
    height = Cm(25 / img1_aspect_radio)

    prs.slide_width = width
    prs.slide_height = height

    progressbar = tqdm(total=len(image_paths),desc="正在生成ppt文件")
    for index, image_path in enumerate(image_paths):
        # 创建一个新幻灯片，使用第一张幻灯片的布局
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 打开图片并计算大小
        img = cv2.imread(image_path)
        aspect_ratio = img.shape[1] / img.shape[0]

        #设定幻灯片长宽

        # 调整图片大小以适应幻灯片
        width = height = 0
        if img.shape[1] > img.shape[0]:
            width = prs.slide_width
            height = width / aspect_ratio
        else:
            height = prs.slide_height
            width = height * aspect_ratio

        # 添加图片
        left = (prs.slide_width - width) / 2
        top = (prs.slide_height - height) / 2
        pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        progressbar.update(1)
    progressbar.close()
    print("生成ppt完成!")
    # 保存PPT文件
    prs.save(os.path.join(output_ppt_path,'output_ppt.pptx'))
    clean_images_in_path(pictemp_path)  # 清空pictemp文件夹

if __name__ == "__main__":
    main()